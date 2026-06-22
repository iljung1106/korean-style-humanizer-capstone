#!/usr/bin/env python3
"""Build an editable PPTX report for Stage05F metric analysis.

The report uses one slide per metric so each metric section stays on a single
page without DOCX table reflow. Text boxes and tables remain editable; only the
metric distribution plots are inserted as images.
"""

from __future__ import annotations

import json
import math
import sys
from pathlib import Path

import pandas as pd

TMP_PPTX = Path("/tmp/stage05f_pptx_pkg")
if TMP_PPTX.exists():
    sys.path.insert(0, str(TMP_PPTX))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - user-facing dependency error
    raise SystemExit("python-pptx가 필요합니다. 예: python3 -m pip install python-pptx") from exc


ROOT = Path(__file__).resolve().parents[2]
LOCAL_REPORT_DIR = ROOT / "outputs" / "local_reports"
OUTDIR = LOCAL_REPORT_DIR / "stage05f_revised_metric_report"
STATS_CSV = OUTDIR / "stage05f_revised_metric_stats.csv"
PLOT_DIR = (
    LOCAL_REPORT_DIR
    / "stage05f_human_ai_rewrite_report_style_curated_filtered"
    / "metric_page_plots"
)
LEXICON_PATH = ROOT / "data" / "processed" / "anti_slop_lexicon.json"
OUTPPTX = LOCAL_REPORT_DIR / "인간_AI_Rewrite_지표_설명_수정본_시각개선_단일슬라이드.pptx"
WORKPPTX = OUTDIR / "인간_AI_Rewrite_지표_설명_수정본_시각개선_단일슬라이드.pptx"


FONT = "NanumBarunGothic"
BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(17, 24, 39)
MUTED = RGBColor(91, 106, 130)
LINE = RGBColor(205, 213, 224)
SOFT_LINE = RGBColor(226, 232, 240)
LIGHT_BLUE = RGBColor(217, 234, 247)
GREEN = RGBColor(221, 239, 216)
RED = RGBColor(247, 215, 215)
ORANGE = RGBColor(251, 230, 194)
GRAY = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)


def fmt(value: float | int | None, ndigits: int = 4) -> str:
    try:
        value = float(value)
    except Exception:
        return "NA"
    if math.isnan(value) or math.isinf(value):
        return "NA"
    return f"{value:.{ndigits}f}"


def pct(value: float | int | None, ndigits: int = 1) -> str:
    try:
        value = float(value)
    except Exception:
        return "NA"
    if math.isnan(value) or math.isinf(value):
        return "NA"
    return f"{value * 100:.{ndigits}f}%"


def effect_label(delta: float) -> str:
    value = abs(float(delta))
    if value >= 0.474:
        return "large"
    if value >= 0.33:
        return "medium"
    if value >= 0.147:
        return "small"
    return "negligible"


def status_fill(status: str) -> RGBColor:
    if status == "개선":
        return GREEN
    if status == "악화":
        return RED
    return ORANGE


def metric_direction(metric: str) -> str:
    if metric == "pos_4gram_diversity":
        return "높을수록 Human 표본에 가까운 경향"
    if metric == "parenthesis_pair_per_1k_chars":
        return "Human 표본 평균과 가까울수록 적절한 경향"
    return "낮을수록 Human 표본에 가까운 경향"


def metric_method(metric: str) -> tuple[str, str]:
    methods = {
        "comma_per_1k_chars": ("본문 길이에 비해 쉼표가 얼마나 자주 쓰이는지를 측정한다.", "쉼표 수 / 본문 문자 수 × 1,000"),
        "sentence_initial_token_repeat_rate": ("각 문장의 첫 어휘가 반복되는 정도를 측정한다.", "반복 출현한 문장 첫 어휘 수 / 전체 문장 첫 어휘 수"),
        "parenthesis_pair_per_1k_chars": ("본문 길이에 비해 괄호 표현이 얼마나 자주 쓰이는지를 측정한다.", "괄호쌍 수 / 본문 문자 수 × 1,000"),
        "modifier_repeat_burst_mass": ("가까운 위치에서 같은 수식 표현이 반복되는 정도를 측정한다.", "반복 간격이 짧을수록 큰 가중치를 부여해 합산"),
        "modifier_repetition_mass": ("수식 표현 전반의 반복 사용량을 측정한다.", "처음 출현을 제외한 반복 수식 표현 수 / 전체 수식 표현 수"),
        "simile_marker_per_1k_chars": ("직유를 나타내는 표현의 사용 밀도를 측정한다.", "직유 표지 수 / 본문 문자 수 × 1,000"),
        "content_modifier_repeat_occurrence_rate": ("내용어 중 수식 성격의 표현이 반복 출현하는 비율을 측정한다.", "반복 출현 수식 표현량 / 전체 수식 표현량"),
        "simile_sentence_rate": ("직유 표현이 포함된 문장의 비율을 측정한다.", "직유 표지 포함 문장 수 / 전체 문장 수"),
        "pos_3gram_repeat_rate": ("세 개 연속 품사 패턴이 반복되는 정도를 측정한다.", "반복 출현한 3-gram 품사 패턴 수 / 전체 3-gram 수"),
        "anti_slop_density": ("AI 생성 문장에서 과대표집된 상투 표현의 밀도를 측정한다.", "상투 표현 사전의 가중 hit / 전체 feature 수"),
        "pos_5gram_repeat_rate": ("다섯 개 연속 품사 패턴이 반복되는 정도를 측정한다.", "반복 출현한 5-gram 품사 패턴 수 / 전체 5-gram 수"),
        "pos_4gram_diversity": ("네 개 연속 품사 패턴의 다양성을 측정한다.", "서로 다른 4-gram 품사 패턴 수 / 전체 4-gram 수"),
        "sentence_length_iqr_ratio": ("문장 길이의 중간 50% 분산 정도를 측정한다.", "문장 길이 사분위 범위 / 문장 길이 중앙값"),
        "sentence_length_cv": ("문장 길이 변동성을 측정한다.", "문장 길이 표준편차 / 평균 문장 길이"),
        "sentence_final_token_repeat_rate": ("문장 끝 표현이 반복되는 정도를 측정한다.", "반복 출현한 문장 마지막 어휘 수 / 전체 문장 마지막 어휘 수"),
    }
    return methods[metric]


def metric_interpretation(metric: str, gap_closure: float) -> tuple[str, str]:
    specific = {
        "anti_slop_density": ("Rewrite는 AI 원문보다 낮지만 Human과의 잔여 차이가 크다.", "사전 기반 지표이므로 장르 관용 표현이 일부 포함될 수 있다."),
        "sentence_length_iqr_ratio": ("Rewrite는 Human 평균에서 더 멀어졌다.", "짧은 문장과 대화문 비율에 민감하므로 CV와 함께 해석한다."),
        "sentence_length_cv": ("Rewrite는 Human 평균에서 더 멀어졌다.", "장면 유형과 대화문 비율에 따라 값이 달라질 수 있다."),
        "sentence_final_token_repeat_rate": ("Rewrite는 문장 끝 표현 반복에서 악화되었다.", "마지막 어휘 기준 지표이므로 문법적 종결 어미만을 분리한 값은 아니다."),
        "pos_4gram_diversity": ("Rewrite는 약한 개선을 보인다.", "표본 길이가 짧으면 다양도 값이 불안정할 수 있다."),
        "parenthesis_pair_per_1k_chars": ("Rewrite는 평균 기준으로 개선됐지만 보조 지표로 해석한다.", "괄호 사용은 작품 형식과 대사 표기 관습의 영향을 받는다."),
    }
    if metric in specific:
        return specific[metric]
    if gap_closure > 0.3:
        return ("Rewrite 평균은 AI 원문보다 Human 평균에 가까워졌다.", "개별 표본의 개선 여부는 paired 개선율을 함께 확인한다.")
    if gap_closure > 0:
        return ("Rewrite 평균은 개선 방향이지만 개선폭은 제한적이다.", "신뢰구간이 넓은 경우 평균 개선 해석에 주의한다.")
    return ("Rewrite 평균은 AI 원문보다 Human 평균에서 더 멀어졌다.", "지표 특성과 표본 분포를 함께 확인해야 한다.")


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_text_frame(tf, font_size: float, bold: bool = False, color: RGBColor = DARK) -> None:
    for paragraph in tf.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    set_text_frame(tf, size, bold, color)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
    return box


def add_box(slide, x, y, w, h, title, body, fill=WHITE, body_size=9.1):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = SOFT_LINE
    shape.line.width = Pt(0.5)
    add_text(slide, x + 0.08, y + 0.07, w - 0.16, 0.22, title, size=8.2, bold=True, color=BLUE)
    add_text(slide, x + 0.08, y + 0.31, w - 0.16, h - 0.36, body, size=body_size, color=DARK)


def add_title(slide, title, subtitle=""):
    add_text(slide, 0.45, 0.28, 9.6, 0.45, title, size=23, bold=True)
    if subtitle:
        add_text(slide, 0.47, 0.78, 9.6, 0.22, subtitle, size=8.8, color=MUTED)


def add_kpi_row(slide, items, x=0.38, y=1.12, w=10.95, h=0.56):
    cell_w = w / len(items)
    for i, (label, value, fill) in enumerate(items):
        x0 = x + i * cell_w
        shape = slide.shapes.add_shape(1, Inches(x0), Inches(y), Inches(cell_w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = SOFT_LINE
        shape.line.width = Pt(0.4)
        add_text(slide, x0 + 0.03, y + 0.05, cell_w - 0.06, 0.16, label, size=6.4, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, x0 + 0.03, y + 0.25, cell_w - 0.06, 0.24, value, size=10.4, bold=True, align=PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, headers, rows, font_size=6.8):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = FONT
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = DARK
                if cell in table.rows[0].cells:
                    paragraph.font.bold = True
    return table_shape


def find_plot(metric: str) -> Path | None:
    candidates = sorted(PLOT_DIR.glob(f"*_{metric}.png"))
    return candidates[0] if candidates else None


def add_cover(prs: Presentation, stats: pd.DataFrame) -> None:
    slide = blank_slide(prs)
    add_title(slide, "Human-AI-Rewrite 문체 지표 보고서", "Stage05F rewrite 평가")
    add_kpi_row(
        slide,
        [
            ("지표", f"{len(stats)}개", LIGHT_BLUE),
            ("개선", f"{int((stats.status == '개선').sum())}개", GREEN),
            ("악화", f"{int((stats.status == '악화').sum())}개", RED),
            ("변화 작음", f"{int((stats.status == '변화 작음').sum())}개", ORANGE),
            ("Human / AI / Rewrite", "300 / 300 / 100", GRAY),
        ],
        y=1.22,
    )
    add_box(slide, 0.45, 2.0, 10.9, 0.95, "보고서 목적", "AI 원문을 rewrite한 결과가 인간 웹소설 표본의 문체 분포에 가까워졌는지 정량 지표로 비교한다. 각 지표는 평균 차이, paired 개선율, 효과크기, 분포 요약을 함께 제시한다.", LIGHT_BLUE, body_size=10.2)
    add_box(slide, 0.45, 3.25, 10.9, 0.95, "해석 원칙", "지표는 문체적 경향을 요약하는 보조 척도이다. 인간 평균에 가까워진 결과가 곧 작품성 전체의 향상을 의미하지는 않으며, 개별 텍스트 검토와 함께 해석해야 한다.", WHITE, body_size=10.2)
    add_box(slide, 0.45, 4.5, 10.9, 0.95, "주요 결과", "쉼표 밀도, 문장 시작 반복률, 수식어 반복, 직유 표지 계열은 개선됐다. 반면 anti-slop density의 잔여 격차와 문장 길이 변동성, 문장 끝 표현 반복은 후속 개선이 필요한 축으로 나타났다.", WHITE, body_size=10.2)


def add_method_slides(prs: Presentation, stats: pd.DataFrame) -> None:
    slide = blank_slide(prs)
    add_title(slide, "실험 개요 및 분석 방법")
    add_box(slide, 0.45, 1.25, 5.1, 1.0, "데이터 구성", "Human 표본 300개, AI 원문 표본 300개, Rewrite 표본 100개를 비교했다. Rewrite는 AI 원문 일부를 같은 과업 조건에서 다시 작성한 결과이다.", LIGHT_BLUE)
    add_box(slide, 0.45, 2.55, 5.1, 1.0, "전처리", "분석 대상 본문을 정규화한 뒤 문장, 어휘, 품사 단위 지표를 산출했다. 길이 차이를 줄이기 위해 밀도형 지표는 문자 수 또는 feature 수로 정규화했다.")
    add_box(slide, 0.45, 3.85, 5.1, 1.0, "Paired 분석", "Rewrite와 대응되는 AI 원문을 짝지어, rewrite 이후 Human 평균과의 거리가 감소했는지 계산했다. paired 개선율은 이 거리 감소가 관측된 비율이다.", GREEN)
    add_box(slide, 6.05, 1.25, 5.1, 1.0, "Gap closure", "gap closure = (|AI 평균 - Human 평균| - |Rewrite 평균 - Human 평균|) / |AI 평균 - Human 평균|. 양수는 개선, 음수는 악화를 뜻한다.", LIGHT_BLUE)
    add_box(slide, 6.05, 2.55, 5.1, 1.0, "Bootstrap CI", "표본 재추출을 통해 gap closure의 95% 신뢰구간을 추정했다. 구간이 넓을수록 평균 개선 해석은 보수적으로 본다.")
    add_box(slide, 6.05, 3.85, 5.1, 1.12, "Cliff's delta", "두 분포의 분리 정도를 나타내는 비모수 효과크기이다. Rewrite vs AI는 rewrite가 AI 원문 분포에서 얼마나 이동했는지, Rewrite vs Human은 인간 표본과의 잔여 차이를 나타낸다.")

    slide = blank_slide(prs)
    add_title(slide, "Anti-slop Lexicon 및 전체 요약")
    lexicon = json.loads(LEXICON_PATH.read_text(encoding="utf-8"))
    terms = lexicon.get("terms", [])
    counts = pd.Series([term.get("kind") for term in terms]).value_counts().to_dict()
    add_kpi_row(
        slide,
        [
            ("lexicon term", f"{len(terms)}개", LIGHT_BLUE),
            ("unigram", str(counts.get("unigram", 0)), GRAY),
            ("bigram", str(counts.get("bigram", 0)), GRAY),
            ("skip-bigram", str(counts.get("skip_bigram", 0)), GRAY),
            ("skip-trigram", str(counts.get("skip_trigram", 0)), GRAY),
        ],
        y=1.05,
    )
    rows = [
        [r.metric_ko, r.status, pct(r.gap_closure), pct(r.paired_improved_rate), fmt(r.rewrite_vs_ai_delta, 3), fmt(r.rewrite_vs_human_delta, 3)]
        for r in stats.itertuples()
    ]
    add_table(slide, 0.45, 1.95, 10.9, 4.1, ["지표", "상태", "gap closure", "paired 개선율", "Rewrite vs AI δ", "Rewrite vs Human δ"], rows, font_size=5.9)


def add_metric_slide(prs: Presentation, index: int, row) -> None:
    slide = blank_slide(prs)
    metric = row.metric
    add_title(slide, f"{index:02d}. {row.metric_ko}", metric)
    add_kpi_row(
        slide,
        [
            ("상태", row.status, status_fill(row.status)),
            ("gap closure", pct(row.gap_closure), LIGHT_BLUE),
            ("bootstrap CI", f"{pct(row.gap_ci_low)}~{pct(row.gap_ci_high)}", GRAY),
            ("paired 개선율", pct(row.paired_improved_rate), GREEN if row.paired_improved_rate >= 0.5 else ORANGE),
            ("Rewrite vs AI δ", fmt(row.rewrite_vs_ai_delta, 3), GRAY),
            ("Rewrite vs Human δ", fmt(row.rewrite_vs_human_delta, 3), GRAY),
        ],
        y=1.08,
    )
    plot = find_plot(metric)
    if plot:
        slide.shapes.add_picture(str(plot), Inches(0.58), Inches(1.82), width=Inches(5.28))
    add_table(
        slide,
        0.45,
        5.45,
        5.62,
        0.86,
        ["group", "mean", "median", "q25", "q75", "std", "n"],
        [
            ["Human", fmt(row.human_mean), fmt(row.human_median), fmt(row.human_q25), fmt(row.human_q75), fmt(row.human_std), "300"],
            ["AI", fmt(row.ai_mean), fmt(row.ai_median), fmt(row.ai_q25), fmt(row.ai_q75), fmt(row.ai_std), "300"],
            ["Rewrite", fmt(row.rewrite_mean), fmt(row.rewrite_median), fmt(row.rewrite_q25), fmt(row.rewrite_q75), fmt(row.rewrite_std), "100"],
        ],
        font_size=5.7,
    )
    prep, calc = metric_method(metric)
    interp, caution = metric_interpretation(metric, row.gap_closure)
    add_box(slide, 6.18, 1.82, 5.02, 0.62, "준비 방식", prep, LIGHT_BLUE, body_size=7.8)
    add_box(slide, 6.18, 2.58, 5.02, 0.62, "계산 방식", calc, WHITE, body_size=7.8)
    add_box(slide, 6.18, 3.34, 5.02, 0.55, "목표 방향", metric_direction(metric), WHITE, body_size=7.8)
    add_box(slide, 6.18, 4.03, 5.02, 0.72, "해석", interp, status_fill(row.status), body_size=7.8)
    result = (
        f"Human={fmt(row.human_mean)}, AI={fmt(row.ai_mean)}, Rewrite={fmt(row.rewrite_mean)}. "
        f"AI vs Human δ={fmt(row.ai_vs_human_delta, 3)} ({effect_label(row.ai_vs_human_delta)}), paired n={int(row.paired_n)}."
    )
    add_box(slide, 6.18, 4.89, 5.02, 0.72, "이번 결과", result, WHITE, body_size=7.5)
    add_box(slide, 6.18, 5.75, 5.02, 0.58, "주의점", caution, WHITE, body_size=7.5)


def build_pptx(stats: pd.DataFrame) -> None:
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)
    add_cover(prs, stats)
    add_method_slides(prs, stats)
    for index, row in enumerate(stats.itertuples(), start=1):
        add_metric_slide(prs, index, row)
    prs.save(WORKPPTX)
    prs.save(OUTPPTX)


def main() -> None:
    stats = pd.read_csv(STATS_CSV).sort_values("gap_closure", ascending=False)
    build_pptx(stats)
    print(json.dumps({"pptx": str(OUTPPTX), "workpptx": str(WORKPPTX), "slides": 3 + len(stats)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
